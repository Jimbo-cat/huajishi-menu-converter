#!/usr/bin/env python3
"""Extract data from 5月 report PPTX files (try alternatives)."""

from pptx import Presentation
import os

base = "/Users/zhengbo/Desktop/月报/1A 2026"
candidates = ["桌面5月工作总结.pptx", "桌面 5 月工作总结.pptx", "桌面 5 月工作总结1.pptx"]

for c in candidates:
    path = os.path.join(base, c)
    if os.path.exists(path):
        print(f"Trying: {c}")
        print(f"File size: {os.path.getsize(path)} bytes")
        try:
            prs = Presentation(path)
            print(f"OK! Slides: {len(prs.slides)}")
            for i, slide in enumerate(prs.slides):
                print(f"\n--- 第{i+1}页 ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            print(text[:200])
                    elif shape.has_table:
                        table = shape.table
                        print(f"[表格: {len(table.rows)}行 x {len(table.columns)}列]")
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip()[:40] for cell in row.cells]
                            print(f"  行{row_idx}: {' | '.join(cells)}")
            break
        except Exception as e:
            print(f"Error: {e}")
            print()
    else:
        print(f"Not found: {c}")
