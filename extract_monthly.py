#!/usr/bin/env python3
"""Extract data from all 2026 monthly report PPTX files."""

from pptx import Presentation
import os

base = "/Users/zhengbo/Desktop/月报/1A 2026"
files = [
    "桌面1月工作总结.pptx",
    "桌面2月工作总结.pptx",
    "桌面3月工作总结.pptx",
    "桌面4月工作总结.pptx",
    "桌面5月工作总结.pptx"
]

for fname in files:
    path = os.path.join(base, fname)
    print(f"\n{'='*60}")
    print(f"=== {fname} ===")
    print(f"{'='*60}")
    prs = Presentation(path)
    print(f"总幻灯片数: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- 第{i+1}页 ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(text[:200])
            elif shape.has_table:
                table = shape.table
                print(f"[表格: {len(table.rows)}行 x {len(table.columns)}列]")
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip()[:40] for cell in row.cells]
                    print(f"  行{row_idx}: {' | '.join(cells)}")
