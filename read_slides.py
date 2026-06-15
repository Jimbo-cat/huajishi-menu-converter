#!/usr/bin/env python3
"""Extract last 4 slides from 4月 and 5月 PPTs for work order data."""
import sys
from pptx import Presentation

for fname in ['桌面4月工作总结.pptx', '桌面 5 月工作总结.pptx']:
    path = f'/Users/zhengbo/Desktop/月报/1A 2026/{fname}'
    prs = Presentation(path)
    print(f"\n=== {fname} (slides 10-13) ===")
    for i, slide in enumerate(prs.slides):
        if i >= 9:
            print(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    print(text[:800])
