#!/usr/bin/env python3
"""
花吉食菜单转换器
================
功能：读取花吉食菜单 Excel 文件，基于模板 PPTX 生成 5 页（周一~周五）菜单 PPTX。
"""

import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import openpyxl
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from datetime import datetime, timedelta

# ── 参数配置 ──────────────────────────────────────────────────────────────
TEMPLATE_PATH = os.path.expanduser("~/Desktop/餐饮花吉食菜单/6月10菜单模板.pptx")
DEFAULT_OUTPUT_DIR = os.path.expanduser("~/Desktop/餐饮花吉食菜单/输出")

# 星期 → Excel 列字母映射（B=周一, D=周二, F=周三, H=周四, J=周五）
DAY_COLUMNS = ['B', 'D', 'F', 'H', 'J']
DAY_NAMES = ['星期一', '星期二', '星期三', '星期四', '星期五']

# PPTX 中需替换内容的文本框名称（所有 Slide 通用）
# ── 早餐区 ──
BREAKFAST_BOXES = {
    '干点1': '文本框 19',     # 鲜肉大包 香菇菜包 等
    '干点2': '文本框 56',     # 小笼包 派鲜包 等（位置偏右）
    '湿点':  '文本框 54',     # 粥/豆浆/面条
    '煎炸':  '文本框 52',     # 煎饼类
    '蛋类':  '文本框 51',     # 蛋类
    '西点':  '文本框 43',     # 面包/蛋糕
    '粗粮':  '文本框 9',      # 玉米/红薯
}
# ── 午餐区 ──
LUNCH_BOXES = {
    '主荤': '文本框 35',
    '副荤': '文本框 44',
    '蔬菜': '文本框 45',
    '甜汤': '文本框 48',
    '例汤': '文本框 49',
    '面档': '文本框 47',
}
# ── 晚餐区 ──
DINNER_BOXES = {
    '主荤': '文本框 32',
    '副荤': '文本框 33',
    '蔬菜': '文本框 34',
    '例汤': '文本框 64',
}


def clean_item(s):
    """清理数据中的残留标记"""
    s = s.strip()
    s = s.replace('/pig', '').replace('（新）', '').strip()
    return s


def get_cell_str(ws, col_letter, row):
    """安全读取单元格，返回字符串（去 None）"""
    v = ws[f'{col_letter}{row}'].value
    return str(v).strip() if v is not None else ''


def get_cell_list(ws, col_letter, start_row, end_row):
    """读取一列中某范围单元格，返回非空字符串列表，自动拆分 \\n 并清理 /pig"""
    items = []
    for r in range(start_row, end_row + 1):
        v = ws[f'{col_letter}{r}'].value
        if v is not None:
            s = str(v).strip()
            if s and s != 'None':
                # 拆分单元格内的换行
                for part in s.split(chr(10)):
                    part = part.strip()
                    if part:
                        items.append(clean_item(part))
    return items


def extract_day_data(ws_breakfast, ws_lunch, day_index):
    """
    提取指定天的所有菜单数据。
    day_index: 0=Mon ... 4=Fri
    返回 (breakfast_data, lunch_data, dinner_data) 三个 dict。
    """
    col = DAY_COLUMNS[day_index]

    # ═══ 早餐（早餐 sheet） ═══
    breakfast = {}
    # 面档 → R3
    breakfast['面档'] = get_cell_list(ws_breakfast, col, 3, 3)
    # 干点 → R4-R12
    breakfast['干点'] = get_cell_list(ws_breakfast, col, 4, 12)
    # 粗粮 → R13-R15
    breakfast['粗粮'] = get_cell_list(ws_breakfast, col, 13, 15)
    # 蛋类 → R16-R18
    breakfast['蛋类'] = get_cell_list(ws_breakfast, col, 16, 18)
    # 煎炸 → R19-R20
    breakfast['煎炸'] = get_cell_list(ws_breakfast, col, 19, 20)
    # 湿点 → R21-R24
    breakfast['湿点'] = get_cell_list(ws_breakfast, col, 21, 24)
    # 西点 → R25-R26
    breakfast['西点'] = get_cell_list(ws_breakfast, col, 25, 26)

    # ═══ 午餐（正餐 sheet R3-R23） ═══
    lunch = {}
    lunch['主荤'] = get_cell_list(ws_lunch, col, 3, 6)       # 大荤 ×4
    lunch['副荤'] = get_cell_list(ws_lunch, col, 7, 10)      # 小荤 ×4
    lunch['蔬菜'] = get_cell_list(ws_lunch, col, 11, 13)     # 蔬菜 ×3
    lunch['甜汤'] = get_cell_list(ws_lunch, col, 14, 14)     # 甜汤
    lunch['咸汤'] = get_cell_list(ws_lunch, col, 15, 15)     # 咸汤 → 例汤
    # 面档 R17-R23（非连续行，分别读取）
    lunch['面档'] = get_cell_list(ws_lunch, col, 17, 23)

    # ═══ 晚餐（正餐 sheet R25-R34） ═══
    dinner = {}
    dinner['主荤'] = get_cell_list(ws_lunch, col, 25, 27)    # 大荤 ×3
    dinner['副荤'] = get_cell_list(ws_lunch, col, 28, 30)    # 小荤 ×3
    dinner['蔬菜'] = get_cell_list(ws_lunch, col, 31, 32)    # 蔬菜 ×2
    dinner['汤']   = get_cell_list(ws_lunch, col, 33, 33)    # 汤

    return breakfast, lunch, dinner


def _detect_canonical_style(text_frame):
    """从文本框检测统一的基准样式（run → paragraph → shape 逐级回退）"""
    from pptx.oxml.ns import qn

    style = {'font_name': None, 'ea_font': None,
             'size': None, 'bold': None, 'color_rgb': None}

    # 1. 扫描所有非空 run，找最完整的样式
    for para in text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if t and t != '(empty)':
                s = _read_run_style(run)
                # 优先取既有 font_name 又有 ea_font 的，且尽量有显式颜色
                if s['font_name'] and s['ea_font'] and s['color_rgb']:
                    style = s
                    return style  # 最完整：字体+东亚+颜色全有
                if s['font_name'] and s['ea_font']:
                    # 如果首次遇到带完整字体的，先记下来
                    if not (style.get('font_name') and style.get('ea_font')):
                        style = {**style, **{k: v for k, v in s.items() if v is not None}}
                elif s['font_name'] or s['ea_font']:
                    if not style.get('font_name') and not style.get('ea_font'):
                        style = {**style, **{k: v for k, v in s.items() if v is not None}}

    # 2. run 级别不够 → 检查段落默认样式 (defRPr)
    for para in text_frame.paragraphs:
        pPr = para._p.find(qn('a:pPr'))
        if pPr is not None:
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is not None:
                if not style['font_name']:
                    style['font_name'] = defRPr.get('sz', {}).get('typeface')
                    # Actually sz is the size attribute, let me check
                if not style['ea_font']:
                    ea = defRPr.find(qn('a:ea'))
                    if ea is not None:
                        style['ea_font'] = ea.get('typeface')
                if not style['size']:
                    sz = defRPr.get('sz')
                    if sz:
                        try:
                            from pptx.util import Pt
                            style['size'] = Pt(int(sz) / 100)
                        except:
                            pass
                if not style['bold']:
                    b = defRPr.get('b')
                    if b:
                        style['bold'] = b == '1' or b == 'true'

    # 3. 如果还没有，检查父级 txBody 下的默认段落样式
    if not style['font_name'] and not style['ea_font']:
        try:
            txBody = text_frame._txBody
            for lvl in range(1, 10):
                lvl_pPr = txBody.find(qn(f'a:lvl{lvl}pPr'))
                if lvl_pPr is not None:
                    defRPr = lvl_pPr.find(qn('a:defRPr'))
                    if defRPr is not None:
                        if not style['ea_font']:
                            ea = defRPr.find(qn('a:ea'))
                            if ea is not None:
                                style['ea_font'] = ea.get('typeface')
                        break
        except:
            pass

    return style


def set_textbox_content(text_frame, items, line_spacing_pct=None, color_split=None):
    """
    替换文本框内容，保留原有字体/字号/颜色/加粗/东亚字体。
    每个段落保留其自身的原始样式（面档保持黑色，湿点保持紫色）。
    items: 字符串列表，每个元素成为新的一行（段落）。
    line_spacing_pct: 可选，OpenXML 行距百分比值，如 100000=100%。
    color_split: 可选，[(count, style), ...] 按段次应用不同颜色样式。
               例：[(3, 黑色样式), (2, 紫色样式)] 前3项用样式1，后2项用样式2。
    """
    # 先读取每个段落的原始样式（只读有内容的段落）
    orig_styles = []
    empty_style = None
    for para in text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            # 只从有内容的 run 读样式
            if t:
                s = _read_run_style(run)
                orig_styles.append(s)
                break
        else:
            # 段落有 run 但全空，用第一个
            if para.runs:
                s = _read_run_style(para.runs[0])
                orig_styles.append(s)
                if empty_style is None:
                    empty_style = s
            else:
                orig_styles.append(None)

    if empty_style is None:
        # 尝试从任何 run 找一个
        for para in text_frame.paragraphs:
            for run in para.runs:
                s = _read_run_style(run)
                if s.get('font_name') or s.get('ea_font'):
                    empty_style = s
                    break
            if empty_style:
                break
        if empty_style is None and text_frame.paragraphs and text_frame.paragraphs[0].runs:
            empty_style = _read_run_style(text_frame.paragraphs[0].runs[0])

    if not empty_style:
        empty_style = {'font_name': None, 'ea_font': None,
                       'size': None, 'bold': None, 'color_rgb': None}

    # 备份：如果模板没有足够多段落样式，新段落用最后一个
    if not orig_styles:
        orig_styles = [empty_style]

    # 清空所有段落，只保留第一个 run
    for para in text_frame.paragraphs:
        runs = para.runs
        if runs:
            for r in list(runs)[1:]:
                r._r.getparent().remove(r._r)
            runs[0].text = ''
        else:
            para.add_run()

    # 写入内容
    for i, item in enumerate(items):
        if i < len(text_frame.paragraphs):
            para = text_frame.paragraphs[i]
            runs = para.runs
            if len(runs) > 1:
                for r in list(runs)[1:]:
                    r._r.getparent().remove(r._r)
            if runs:
                run = runs[0]
                run.text = ''
            else:
                run = para.add_run()
        else:
            para = text_frame.add_paragraph()
            if not para.runs:
                para.add_run()
            run = para.runs[0]
        # 写文本
        run.text = str(item)

        # 使用 color_split 或逐段样式
        if color_split:
            offset = 0
            selected_style = empty_style
            for count, split_style in color_split:
                if i < offset + count:
                    selected_style = split_style
                    break
                offset += count
            _apply_run_style(run, selected_style or empty_style)
        else:
            # 用对应段落的原始样式（越界用最后一个）
            style_idx = min(i, len(orig_styles) - 1)
            style = orig_styles[style_idx] if orig_styles else empty_style
            _apply_run_style(run, style or empty_style)

        # 应用行距
        if line_spacing_pct is not None:
            from pptx.oxml.ns import qn
            pPr = para._p.find(qn('a:pPr'))
            if pPr is None:
                pPr = para._p.makeelement(qn('a:pPr'), {})
                para._p.insert(0, pPr)
            # 清除现有行距设置
            old_lnSpc = pPr.find(qn('a:lnSpc'))
            if old_lnSpc is not None:
                pPr.remove(old_lnSpc)
            lnSpc = para._p.makeelement(qn('a:lnSpc'), {})
            spcPct = para._p.makeelement(qn('a:spcPct'), {'val': str(line_spacing_pct)})
            lnSpc.append(spcPct)
            pPr.append(lnSpc)

    # 隐藏多余段落
    for j in range(len(items), len(text_frame.paragraphs)):
        p = text_frame.paragraphs[j]
        runs = p.runs
        if runs:
            for r in list(runs)[1:]:
                r._r.getparent().remove(r._r)
            runs[0].text = ''
        if not p.runs:
            p.add_run()


def set_textbox_two_column(text_frame, items, line_spacing_pct=None, color_split=None):
    """
    两列布局版 set_textbox_content。
    将 items 两两分组，每组占一个段落，段落内两个 run 分别写两个菜品。
    奇数项则最后一个段落只有一个 run。
    color_split: 可选，[(count, style), ...] 逐项（不是逐段）应用样式。
    """
    # 先读取模板的段落样式（只读有内容的段落）
    orig_styles = []
    empty_style = None
    for para in text_frame.paragraphs:
        first_run = para.runs[0] if para.runs else None
        if first_run:
            s = _read_run_style(first_run)
            orig_styles.append(s)
            if empty_style is None and (s.get('font_name') or s.get('ea_font')):
                empty_style = s
        else:
            orig_styles.append(None)

    if empty_style is None:
        for para in text_frame.paragraphs:
            for run in para.runs:
                s = _read_run_style(run)
                if s.get('font_name') or s.get('ea_font'):
                    empty_style = s
                    break
            if empty_style:
                break
        if empty_style is None and text_frame.paragraphs and text_frame.paragraphs[0].runs:
            empty_style = _read_run_style(text_frame.paragraphs[0].runs[0])

    if not empty_style:
        empty_style = {'font_name': None, 'ea_font': None,
                       'size': None, 'bold': None, 'color_rgb': None, 'scheme_clr': None}
    if not orig_styles:
        orig_styles = [empty_style]

    # 清空所有段落，重置
    for para in text_frame.paragraphs:
        runs = para.runs
        if runs:
            for r in list(runs)[1:]:
                r._r.getparent().remove(r._r)
            runs[0].text = ''
        else:
            para.add_run()

    # 将 items 两两分组
    pairs = []
    for i in range(0, len(items), 2):
        pair = [items[i]]
        if i + 1 < len(items):
            pair.append(items[i + 1])
        pairs.append(pair)

    # 写入数据
    for i, pair in enumerate(pairs):
        if i < len(text_frame.paragraphs):
            para = text_frame.paragraphs[i]
            runs = para.runs
            # 确保有足够的 run
            needed = len(pair)
            current = len(runs)
            if current > needed:
                for r in list(runs)[needed:]:
                    r._r.getparent().remove(r._r)
            elif current < needed:
                for _ in range(needed - current):
                    para.add_run()
            # 清空已有文本
            for r in para.runs:
                r.text = ''
            # 写文本
            for j, txt in enumerate(pair):
                para.runs[j].text = str(txt)
        else:
            para = text_frame.add_paragraph()
            for txt in pair:
                r = para.add_run()
                r.text = str(txt)

        # 应用样式到每个 run
        for j in range(len(pair)):
            run = para.runs[j]
            # 确定使用哪个样式
            item_idx = i * 2 + j  # 原始 items 中的索引
            if color_split:
                offset = 0
                selected_style = empty_style
                for count, split_style in color_split:
                    if item_idx < offset + count:
                        selected_style = split_style
                        break
                    offset += count
                _apply_run_style(run, selected_style or empty_style)
            else:
                style_idx = min(item_idx, len(orig_styles) - 1)
                style = orig_styles[style_idx] if orig_styles else empty_style
                _apply_run_style(run, style or empty_style)

        # 行距
        if line_spacing_pct is not None:
            from pptx.oxml.ns import qn
            pPr = para._p.find(qn('a:pPr'))
            if pPr is None:
                pPr = para._p.makeelement(qn('a:pPr'), {})
                para._p.insert(0, pPr)
            old_lnSpc = pPr.find(qn('a:lnSpc'))
            if old_lnSpc is not None:
                pPr.remove(old_lnSpc)
            lnSpc = para._p.makeelement(qn('a:lnSpc'), {})
            spcPct = para._p.makeelement(qn('a:spcPct'), {'val': str(line_spacing_pct)})
            lnSpc.append(spcPct)
            pPr.append(lnSpc)

    # 隐藏多余段落
    for j in range(len(pairs), len(text_frame.paragraphs)):
        p = text_frame.paragraphs[j]
        runs = p.runs
        if runs:
            for r in list(runs)[1:]:
                r._r.getparent().remove(r._r)
            runs[0].text = ''
        if not p.runs:
            p.add_run()


def _read_run_style(run):
    """读取 run 的完整样式（含 scheme 主题色）"""
    from pptx.oxml.ns import qn
    from pptx.dml.color import RGBColor

    style = {}
    style['font_name'] = run.font.name  # Latin font

    # 东亚字体（从 XML 读 ea 标签）
    ea_font = None
    try:
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            ea = rPr.find(qn('a:ea'))
            if ea is not None:
                ea_font = ea.get('typeface')
    except Exception:
        pass
    style['ea_font'] = ea_font

    style['size'] = run.font.size
    style['bold'] = run.font.bold

    # 显式 RGB 颜色
    color_rgb = None
    try:
        c = run.font.color
        if c and c.type is not None:
            color_rgb = str(c.rgb)
    except Exception:
        pass
    style['color_rgb'] = color_rgb

    # 主题色（schemeClr，如 accent2=橙色）
    scheme_clr = None
    try:
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            fill = rPr.find(qn('a:solidFill'))
            if fill is not None:
                scheme = fill.find(qn('a:schemeClr'))
                if scheme is not None:
                    scheme_clr = scheme.get('val')
    except Exception:
        pass
    style['scheme_clr'] = scheme_clr

    # 主题色的亮度修饰（lumMod/lumOff）— 如 tx1 的 lumMod 65000 + lumOff 35000 = 灰色
    lum_mod = None
    lum_off = None
    try:
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            fill = rPr.find(qn('a:solidFill'))
            if fill is not None:
                scheme = fill.find(qn('a:schemeClr'))
                if scheme is not None:
                    lm = scheme.find(qn('a:lumMod'))
                    if lm is not None:
                        lum_mod = lm.get('val')
                    lo = scheme.find(qn('a:lumOff'))
                    if lo is not None:
                        lum_off = lo.get('val')
    except Exception:
        pass
    style['lum_mod'] = lum_mod
    style['lum_off'] = lum_off

    return style


def _apply_run_style(run, style):
    """将样式应用到 run"""
    from pptx.oxml.ns import qn
    from pptx.dml.color import RGBColor

    # Latin 字体
    if style.get('font_name'):
        run.font.name = style['font_name']

    # 东亚字体（通过 XML 设置 ea 标签）
    if style.get('ea_font'):
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                from lxml import etree
                ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', style['ea_font'])
        except Exception:
            pass

    # 字号
    if style.get('size'):
        run.font.size = style['size']

    # 加粗
    run.font.bold = style.get('bold')

    # 颜色 - 按优先顺序：scheme_clr > color_rgb > 不处理
    if style.get('color_rgb'):
        # 显式 RGB 颜色
        try:
            run.font.color.rgb = RGBColor(*bytes.fromhex(style['color_rgb']))
        except Exception:
            pass
    elif style.get('scheme_clr'):
        # 主题色（如 accent2=橙色），写 XML
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            rPr = run._r.get_or_add_rPr()
            # 删除现有 solidFill
            for old_fill in list(rPr.findall(qn('a:solidFill'))):
                rPr.remove(old_fill)
            # 创建 solidFill → schemeClr
            fill = etree.SubElement(rPr, qn('a:solidFill'))
            scheme = etree.SubElement(fill, qn('a:schemeClr'))
            scheme.set('val', style['scheme_clr'])
            # 恢复亮度修饰（如 tx1 的 lumMod 65000 + lumOff 35000 = 灰色）
            if style.get('lum_mod'):
                lm = etree.SubElement(scheme, qn('a:lumMod'))
                lm.set('val', str(style['lum_mod']))
            if style.get('lum_off'):
                lo = etree.SubElement(scheme, qn('a:lumOff'))
                lo.set('val', str(style['lum_off']))
        except Exception:
            pass
    # 没有颜色信息 → 保持 run 原有颜色不变


def _detect_dual_styles(text_frame):
    """
    从文本框检测两种不同的颜色样式。
    返回：(style_nocolor, style_color) — 前者是黑色/继承色样式，后者是显色样式。
    如果只有一种样式，两个值相同。
    """
    style_nocolor = None
    style_color = None
    for para in text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if t and t != '(empty)':
                s = _read_run_style(run)
                # 有显式颜色 vs 无颜色（tx1/黑色不算"有颜色"）
                if s.get('color_rgb'):
                    if style_color is None:
                        style_color = s
                else:
                    if style_nocolor is None:
                        style_nocolor = s
                if style_color and style_nocolor:
                    return (style_nocolor, style_color)
    # 没找到两种：用找到的那个，没有就用 None
    if style_color:
        return (style_color, style_color)
    if style_nocolor:
        return (style_nocolor, style_nocolor)
    return ({'font_name': None, 'ea_font': None, 'size': None, 'bold': None, 'color_rgb': None},
            {'font_name': None, 'ea_font': None, 'size': None, 'bold': None, 'color_rgb': None})


def find_textbox(slide, box_name):
    """根据 name 在 slide 上找到文本框，返回 Shape 或 None"""
    for shape in slide.shapes:
        if shape.name == box_name and shape.has_text_frame:
            return shape
    return None


def generate_pptx(excel_path, output_path, template_path=None):
    """
    核心转换函数：读取 Excel → 生成 PPTX。
    返回 (bool, str) → (成功?, 消息)
    """
    if template_path is None:
        template_path = TEMPLATE_PATH

    if not os.path.exists(excel_path):
        return False, f"Excel 文件不存在: {excel_path}"
    if not os.path.exists(template_path):
        return False, f"模板文件不存在: {template_path}"

    try:
        # 1. 读取 Excel
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        if '早餐' not in wb.sheetnames or '正餐' not in wb.sheetnames:
            return False, "Excel 中需要包含【早餐】和【正餐】两个 Sheet"
        ws_breakfast = wb['早餐']
        ws_lunch = wb['正餐']

        # 2. 打开模板
        prs = Presentation(template_path)

        # 3. 处理每一页（5页 = 周一到周五）
        for slide_idx in range(min(5, len(prs.slides))):
            slide = prs.slides[slide_idx]
            breakfast, lunch, dinner = extract_day_data(ws_breakfast, ws_lunch, slide_idx)

            # ── 更新早餐区 ──
            # 干点（文本框 19）：两列布局（模板原就是两列），取前 5 个干点项
            dry_items = breakfast['干点']
            tb19 = find_textbox(slide, '文本框 19')
            if tb19 and dry_items:
                set_textbox_two_column(tb19.text_frame, dry_items[:5])
            elif tb19:
                set_textbox_content(tb19.text_frame, [''])

            # 干点2（文本框 56）：剩余干点（特色项如小笼包/生煎）
            tb56 = find_textbox(slide, '文本框 56')
            if tb56:
                dry_remaining = dry_items[5:] if len(dry_items) > 5 else []
                set_textbox_content(tb56.text_frame, dry_remaining if dry_remaining else [''])

            # 湿点（文本框 54）：面档 + 湿点合并，两列布局
            # 统一使用紫色，所有项颜色一致
            wet_items = breakfast['面档'] + breakfast['湿点']
            tb54 = find_textbox(slide, '文本框 54')
            if tb54 and wet_items:
                split_styles = _detect_dual_styles(tb54.text_frame)
                # 统一使用紫色样式（split_styles[1]），湿点所有项颜色一致
                uniform_style = split_styles[1]
                color_split = [(len(wet_items), uniform_style)]
                set_textbox_two_column(tb54.text_frame, wet_items, color_split=color_split)
            elif tb54:
                set_textbox_content(tb54.text_frame, [''])

            # 煎炸（文本框 52）
            tb52 = find_textbox(slide, '文本框 52')
            if tb52:
                set_textbox_content(tb52.text_frame, breakfast['煎炸'] if breakfast['煎炸'] else [''])

            # 蛋类（文本框 51）
            tb51 = find_textbox(slide, '文本框 51')
            if tb51:
                set_textbox_content(tb51.text_frame, breakfast['蛋类'] if breakfast['蛋类'] else [''], line_spacing_pct=150000)

            # 西点（文本框 43）
            tb43 = find_textbox(slide, '文本框 43')
            if tb43:
                set_textbox_content(tb43.text_frame, breakfast['西点'] if breakfast['西点'] else [''])

            # 粗粮（文本框 1）
            tb1 = find_textbox(slide, '文本框 1')
            if tb1:
                set_textbox_content(tb1.text_frame, breakfast['粗粮'] if breakfast['粗粮'] else [''])

            # ── 更新午餐区 ──
            # 主荤（文本框 35）
            tb35 = find_textbox(slide, '文本框 35')
            if tb35:
                set_textbox_content(tb35.text_frame, lunch['主荤'] if lunch['主荤'] else [''])

            # 副荤（文本框 44）
            tb44 = find_textbox(slide, '文本框 44')
            if tb44:
                set_textbox_content(tb44.text_frame, lunch['副荤'] if lunch['副荤'] else [''])

            # 蔬菜（文本框 45）
            tb45 = find_textbox(slide, '文本框 45')
            if tb45:
                set_textbox_content(tb45.text_frame, lunch['蔬菜'] if lunch['蔬菜'] else [''])

            # 甜汤（文本框 48）
            tb48 = find_textbox(slide, '文本框 48')
            if tb48:
                set_textbox_content(tb48.text_frame, lunch['甜汤'] if lunch['甜汤'] else [''])

            # 例汤（文本框 49）→ 用咸汤
            tb49 = find_textbox(slide, '文本框 49')
            if tb49:
                set_textbox_content(tb49.text_frame, lunch['咸汤'] if lunch['咸汤'] else [''])

            # 面档（文本框 47）— 两列布局，保留橙色主题色
            tb47 = find_textbox(slide, '文本框 47')
            if tb47 and lunch['面档']:
                set_textbox_two_column(tb47.text_frame, lunch['面档'])
            elif tb47:
                set_textbox_content(tb47.text_frame, [''])

            # ── 更新晚餐区 ──
            # 主荤（文本框 32）
            tb32 = find_textbox(slide, '文本框 32')
            if tb32:
                set_textbox_content(tb32.text_frame, dinner['主荤'] if dinner['主荤'] else [''])

            # 副荤（文本框 33）
            tb33 = find_textbox(slide, '文本框 33')
            if tb33:
                set_textbox_content(tb33.text_frame, dinner['副荤'] if dinner['副荤'] else [''])

            # 蔬菜（文本框 34）
            tb34 = find_textbox(slide, '文本框 34')
            if tb34:
                set_textbox_content(tb34.text_frame, dinner['蔬菜'] if dinner['蔬菜'] else [''])

            # 例汤（文本框 64）
            tb64 = find_textbox(slide, '文本框 64')
            if tb64:
                set_textbox_content(tb64.text_frame, dinner['汤'] if dinner['汤'] else [''])

        # 4. 保存
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        prs.save(output_path)
        return True, f"生成成功！\n输出文件: {output_path}"

    except Exception as e:
        return False, f"生成失败: {str(e)}"


# ══════════════════════════════════════════════════════════════════════════
#  GUI 界面（tkinter）
# ══════════════════════════════════════════════════════════════════════════

class MenuConverterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("花吉食菜单转换器")
        self.root.geometry("800x500")
        self.root.resizable(False, False)

        # 状态变量
        self.excel_path = tk.StringVar()
        self.output_dir = tk.StringVar(value=DEFAULT_OUTPUT_DIR)
        self.template_path = tk.StringVar(value=TEMPLATE_PATH)

        self._build_ui()
        self._center_window()

    def _build_ui(self):
        # ── 主容器 ──
        main_frame = ttk.Frame(self.root, padding=20)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # 标题
        title = ttk.Label(main_frame, text="花吉食菜单转换器",
                          font=('PingFang SC', 18, 'bold'))
        title.pack(pady=(0, 20))

        # ── 左右分栏容器 ──
        content_frame = ttk.Frame(main_frame)
        content_frame.pack(fill=tk.BOTH, expand=True)

        # ── 左侧面板：选择文件 ──
        left_frame = ttk.LabelFrame(content_frame, text="📄 选择菜单模板（Excel）", padding=15)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))

        ttk.Label(left_frame, text="花吉食菜单.xlsx：", font=('PingFang SC', 10)).pack(anchor=tk.W, pady=(0, 5))
        path_frame1 = ttk.Frame(left_frame)
        path_frame1.pack(fill=tk.X, pady=(0, 10))
        self.excel_entry = ttk.Entry(path_frame1, textvariable=self.excel_path, width=35)
        self.excel_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(path_frame1, text="浏览…", command=self._browse_excel).pack(side=tk.RIGHT, padx=(5, 0))

        ttk.Label(left_frame, text="模板 PPTX（可选）：", font=('PingFang SC', 10)).pack(anchor=tk.W, pady=(0, 5))
        path_frame2 = ttk.Frame(left_frame)
        path_frame2.pack(fill=tk.X)
        self.tpl_entry = ttk.Entry(path_frame2, textvariable=self.template_path, width=35)
        self.tpl_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(path_frame2, text="浏览…", command=self._browse_template).pack(side=tk.RIGHT, padx=(5, 0))

        # 提示
        info_box = tk.Text(left_frame, height=12, width=40, font=('PingFang SC', 9),
                           bg='#f5f5f5', relief=tk.FLAT, wrap=tk.WORD, padx=10, pady=10)
        info_box.pack(fill=tk.BOTH, expand=True, pady=(15, 0))
        info_box.insert(tk.END, "使用说明：\n\n")
        info_box.insert(tk.END, "① 点击「浏览…」选择花吉食菜单.xlsx\n\n")
        info_box.insert(tk.END, "② 右侧选择输出目录\n\n")
        info_box.insert(tk.END, "③ 点击右侧「确认生成」按钮\n\n")
        info_box.insert(tk.END, "④ 程序会自动生成 5 页 PPTX\n   （周一至周五）")
        info_box.config(state=tk.DISABLED)

        # ── 中间虚线 ──
        center_frame = ttk.Frame(content_frame, width=30)
        center_frame.pack(side=tk.LEFT, fill=tk.Y, padx=5)
        center_frame.pack_propagate(False)

        # 虚线（画布）
        canvas = tk.Canvas(center_frame, width=2, height=280, highlightthickness=0)
        canvas.pack(expand=True)
        canvas.create_line(1, 0, 1, 280, dash=(4, 4), fill='#999999')

        # ── 右侧面板：输出设置 ──
        right_frame = ttk.LabelFrame(content_frame, text="💾 输出设置", padding=15)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=(10, 0))

        ttk.Label(right_frame, text="输出目录：", font=('PingFang SC', 10)).pack(anchor=tk.W, pady=(0, 5))
        out_frame = ttk.Frame(right_frame)
        out_frame.pack(fill=tk.X, pady=(0, 15))
        self.out_entry = ttk.Entry(out_frame, textvariable=self.output_dir, width=35)
        self.out_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(out_frame, text="浏览…", command=self._browse_output).pack(side=tk.RIGHT, padx=(5, 0))

        # 文件名预览
        ttk.Label(right_frame, text="输出文件名：", font=('PingFang SC', 10)).pack(anchor=tk.W, pady=(0, 5))
        filename = f"花吉食菜单_{(datetime.now().strftime('%Y%m%d'))}.pptx"
        self.filename_label = ttk.Label(right_frame, text=filename, font=('PingFang SC', 9), foreground='#666666')
        self.filename_label.pack(anchor=tk.W, pady=(0, 20))

        # 状态输出
        ttk.Label(right_frame, text="生成状态：", font=('PingFang SC', 10)).pack(anchor=tk.W, pady=(0, 5))
        self.status_text = tk.Text(right_frame, height=5, width=30, font=('PingFang SC', 9),
                                    bg='#f5f5f5', relief=tk.FLAT, wrap=tk.WORD, padx=10, pady=5,
                                    state=tk.DISABLED)
        self.status_text.pack(fill=tk.BOTH, expand=True)

        # 生成按钮（放在右下角）
        btn_frame = ttk.Frame(right_frame)
        btn_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=(10, 0))
        self.generate_btn = ttk.Button(btn_frame, text="确认生成", command=self._generate,
                                       style='Generate.TButton')
        self.generate_btn.pack(side=tk.RIGHT)

        # ── 底部说明 ──
        bottom_label = ttk.Label(main_frame, text="基于模板生成菜单 | 支持 早餐/午餐/晚餐 按分类自动填入",
                                 font=('PingFang SC', 9), foreground='#888888')
        bottom_label.pack(pady=(10, 0))

        # ── 按钮样式 ──
        style = ttk.Style()
        style.configure('Generate.TButton', font=('PingFang SC', 14, 'bold'), padding=(20, 10))

    def _center_window(self):
        self.root.update_idletasks()
        w = self.root.winfo_width()
        h = self.root.winfo_height()
        sw = self.root.winfo_screenwidth()
        sh = self.root.winfo_screenheight()
        x = (sw - w) // 2
        y = (sh - h) // 2
        self.root.geometry(f"{w}x{h}+{x}+{y}")

    def _browse_excel(self):
        path = filedialog.askopenfilename(
            title="选择花吉食菜单 Excel 文件",
            filetypes=[("Excel 文件", "*.xlsx *.xls"), ("所有文件", "*.*")]
        )
        if path:
            self.excel_path.set(path)

    def _browse_template(self):
        path = filedialog.askopenfilename(
            title="选择 PPTX 模板文件",
            filetypes=[("PPTX 模板", "*.pptx"), ("所有文件", "*.*")]
        )
        if path:
            self.template_path.set(path)

    def _browse_output(self):
        path = filedialog.askdirectory(title="选择输出目录")
        if path:
            self.output_dir.set(path)

    def _set_status(self, msg, is_error=False):
        self.status_text.config(state=tk.NORMAL)
        self.status_text.delete(1.0, tk.END)
        if is_error:
            self.status_text.insert(tk.END, msg, 'error')
            self.status_text.tag_config('error', foreground='red')
        else:
            self.status_text.insert(tk.END, msg, 'ok')
            self.status_text.tag_config('ok', foreground='green')
        self.status_text.config(state=tk.DISABLED)

    def _generate(self):
        # 校验
        excel = self.excel_path.get()
        if not excel:
            messagebox.showwarning("提示", "请先选择花吉食菜单 Excel 文件")
            return

        out_dir = self.output_dir.get()
        if not out_dir:
            messagebox.showwarning("提示", "请选择输出目录")
            return

        tpl = self.template_path.get()
        if not tpl:
            tpl = TEMPLATE_PATH

        # 构造输出文件名
        now = datetime.now()
        filename = f"花吉食菜单_{now.strftime('%Y%m%d')}.pptx"
        output_path = os.path.join(out_dir, filename)

        # 执行转换
        self._set_status("正在生成中，请稍候…")
        self.root.update()

        success, msg = generate_pptx(excel, output_path, template_path=tpl)
        self._set_status(msg, is_error=not success)


def main():
    root = tk.Tk()
    app = MenuConverterApp(root)
    root.mainloop()


if __name__ == '__main__':
    main()
